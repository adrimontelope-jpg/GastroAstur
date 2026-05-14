from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colores corporativos
AZUL       = RGBColor(0x1F, 0x4E, 0x79)
AZUL_CLARO = RGBColor(0x2E, 0x75, 0xB6)
TEAL       = RGBColor(0x21, 0x96, 0xA6)
NARANJA    = RGBColor(0xE0, 0x7B, 0x39)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
GRIS       = RGBColor(0x44, 0x44, 0x44)
GRIS_CLARO = RGBColor(0xF0, 0xF4, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=BLANCO, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def blank_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)

# ── DIAPOSITIVA 1: PORTADA ──────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, AZUL)
add_rect(s, 0, 5.5, 13.33, 2.0, TEAL)
add_rect(s, 0, 0, 0.3, 7.5, NARANJA)

add_text(s, "GastroAstur", 0.6, 1.0, 12, 1.5, size=54, bold=True, color=BLANCO)
add_text(s, "Análisis exploratorio de la gastronomía en Asturias",
         0.6, 2.6, 12, 0.8, size=24, color=RGBColor(0xBE, 0xD9, 0xF5))
add_text(s, "Curso de Especialización en IA y Big Data  ·  Convocatoria 2s2526",
         0.6, 5.7, 10, 0.5, size=16, color=BLANCO)
add_text(s, "Adrián Montelpe", 0.6, 6.3, 6, 0.5, size=18, bold=True, color=BLANCO)
add_text(s, "github.com/adrimontelope-jpg/GastroAstur",
         0.6, 6.8, 10, 0.4, size=13, color=RGBColor(0xBE, 0xD9, 0xF5), italic=True)

# ── DIAPOSITIVA 2: ÍNDICE ───────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
add_rect(s, 0, 0, 13.33, 1.2, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, TEAL)

add_text(s, "Índice", 0.6, 0.2, 12, 0.8, size=32, bold=True, color=BLANCO)

items = [
    ("01", "Motivación y contexto"),
    ("02", "Fuentes de datos y metodología"),
    ("03", "Pipeline de ingesta"),
    ("04", "Limpieza y normalización"),
    ("05", "Análisis exploratorio — Hallazgos"),
    ("06", "Modelo de clasificación"),
    ("07", "Dashboard interactivo"),
    ("08", "Conclusiones y vías futuras"),
]
for i, (num, title) in enumerate(items):
    col = i % 2
    row = i // 2
    x = 0.6 + col * 6.4
    y = 1.4 + row * 1.35
    add_rect(s, x, y, 5.9, 1.1, AZUL_CLARO)
    add_text(s, num, x + 0.15, y + 0.05, 0.8, 0.9, size=28, bold=True, color=NARANJA)
    add_text(s, title, x + 0.9, y + 0.2, 4.8, 0.7, size=18, color=BLANCO)

# ── DIAPOSITIVA 3: MOTIVACIÓN ───────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, BLANCO)
add_rect(s, 0, 0, 13.33, 1.2, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, TEAL)

add_text(s, "01 · Motivación y contexto", 0.6, 0.2, 12, 0.8, size=28, bold=True, color=BLANCO)

bullets = [
    "Soy asturiano — la gastronomía es una seña de identidad regional clave.",
    "No existe ninguna herramienta pública de análisis gastronómico para Asturias.",
    "Datos disponibles pero sin explotar: OSM, INE, Wikidata.",
    "Oportunidad de aplicar el ciclo completo de Big Data a un tema real y cercano.",
]
for i, b in enumerate(bullets):
    add_rect(s, 0.6, 1.5 + i*1.35, 0.08, 0.7, NARANJA)
    add_text(s, b, 0.9, 1.45 + i*1.35, 11.8, 0.9, size=20, color=GRIS)

add_rect(s, 0.6, 6.8, 11.8, 0.05, TEAL)
add_text(s, "\"Si hay datos, hay respuestas\"", 0.6, 6.9, 12, 0.4,
         size=14, italic=True, color=AZUL_CLARO, align=PP_ALIGN.CENTER)

# ── DIAPOSITIVA 4: FUENTES Y METODOLOGÍA ────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, BLANCO)
add_rect(s, 0, 0, 13.33, 1.2, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, TEAL)

add_text(s, "02 · Fuentes de datos y metodología", 0.6, 0.2, 12, 0.8, size=28, bold=True, color=BLANCO)

fuentes = [
    ("OpenStreetMap", "Overpass API", "7.568 elementos\nrestaurantes, bares, cafés, pubs"),
    ("INE", "API REST tabla 2886", "Población 2025\n78 municipios asturianos"),
    ("Nominatim", "OSM Geocoding", "Coordenadas oficiales\n15 concejos principales"),
    ("Wikidata", "SPARQL", "Platos y productos\ntípicos asturianos"),
]
for i, (ico, nombre, api, desc) in enumerate(fuentes):
    x = 0.6 + i * 3.1
    add_rect(s, x, 1.4, 2.9, 4.5, GRIS_CLARO)
    add_rect(s, x, 1.4, 2.9, 0.6, AZUL_CLARO)
    add_text(s, ico + " " + nombre, x + 0.1, 1.45, 2.7, 0.5, size=18, bold=True, color=BLANCO)
    add_text(s, api, x + 0.1, 2.1, 2.7, 0.5, size=14, color=AZUL_CLARO, italic=True)
    add_text(s, desc, x + 0.1, 2.7, 2.7, 1.5, size=15, color=GRIS)

add_text(s, "Metodología: ETL → EDA → Modelado → Dashboard  (CRISP-DM)",
         0.6, 6.3, 12, 0.6, size=18, bold=True, color=AZUL, align=PP_ALIGN.CENTER)

# ── DIAPOSITIVA 5: PIPELINE INGESTA ─────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, BLANCO)
add_rect(s, 0, 0, 13.33, 1.2, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, TEAL)

add_text(s, "03 · Pipeline de ingesta", 0.6, 0.2, 12, 0.8, size=28, bold=True, color=BLANCO)

pasos = ["Overpass API", "INE API", "Nominatim", "Wikidata", "data/raw/ JSON"]
cols  = [AZUL_CLARO, TEAL, NARANJA, AZUL_CLARO, RGBColor(0x27,0xAE,0x60)]
for i, (paso, col) in enumerate(zip(pasos, cols)):
    x = 0.5 + i * 2.45
    add_rect(s, x, 2.0, 2.1, 1.2, col)
    add_text(s, paso, x + 0.05, 2.2, 2.0, 0.8, size=16, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(s, "→", x + 2.1, 2.3, 0.35, 0.6, size=28, bold=True, color=AZUL)

add_text(s, "src/ingesta.py", 5.5, 3.6, 2.5, 0.5, size=14, color=AZUL_CLARO, italic=True, align=PP_ALIGN.CENTER)

detalles = [
    "✓  Pausa de 30s entre peticiones (rate limiting)",
    "✓  Metadatos de descarga (fecha, fuente, total elementos)",
    "✓  Manejo de errores 429 y 504 con try/except",
    "✓  Almacenamiento en JSON con ensure_ascii=False",
]
for i, d in enumerate(detalles):
    add_text(s, d, 0.8, 4.3 + i*0.55, 12, 0.5, size=16, color=GRIS)

# ── DIAPOSITIVA 6: LIMPIEZA ──────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, BLANCO)
add_rect(s, 0, 0, 13.33, 1.2, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, TEAL)

add_text(s, "04 · Limpieza y normalización", 0.6, 0.2, 12, 0.8, size=28, bold=True, color=BLANCO)

add_text(s, "Antes de limpiar: 7.568 elementos", 0.6, 1.35, 5.5, 0.5, size=18, bold=True, color=AZUL)
add_text(s, "Después de limpiar: 1.562 registros útiles", 7.0, 1.35, 6.0, 0.5, size=18, bold=True, color=RGBColor(0x27,0xAE,0x60))

campos = [("lat/lon","4%"), ("amenity","24%"), ("nombre","31%"), ("ciudad","81%"), ("cocina","97%")]
for i, (campo, pct) in enumerate(campos):
    x = 0.6 + i * 2.4
    val = int(pct.replace("%",""))
    add_rect(s, x, 6.2 - val*0.045, 2.0, val*0.045, RGBColor(0xE0,0x7B,0x39) if val > 50 else AZUL_CLARO)
    add_text(s, campo, x, 6.25, 2.0, 0.4, size=14, bold=True, color=GRIS, align=PP_ALIGN.CENTER)
    add_text(s, pct + " nulos", x, 6.6, 2.0, 0.3, size=12, color=GRIS, align=PP_ALIGN.CENTER)

add_text(s, "% de valores nulos por campo", 0.6, 1.9, 12, 0.4, size=14, italic=True, color=GRIS)

acciones = [
    "→  dropna(subset=['nombre']): elimina 6.006 elementos sin nombre",
    "→  Normalización de ciudades: 'gijón/xixón' → 'gijón', 'oviedo/uviéu' → 'oviedo'",
    "→  Lowercase y strip en nombre, ciudad y cocina",
    "→  drop_duplicates(subset='osm_id'): elimina duplicados",
]
for i, a in enumerate(acciones):
    add_text(s, a, 0.6, 2.5 + i*0.6, 12.5, 0.55, size=15, color=GRIS)

# ── DIAPOSITIVA 7: HALLAZGOS ─────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, BLANCO)
add_rect(s, 0, 0, 13.33, 1.2, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, TEAL)

add_text(s, "05 · Análisis exploratorio — Hallazgos", 0.6, 0.2, 12, 0.8, size=28, bold=True, color=BLANCO)

kpis = [
    ("7.568", "elementos OSM"),
    ("1.562", "registros útiles"),
    ("7,34", "locales/1000 hab\n(Taramundi)"),
    ("461", "restaurantes\nde cocina regional"),
]
for i, (val, label) in enumerate(kpis):
    x = 0.5 + i * 3.1
    add_rect(s, x, 1.4, 2.9, 1.8, AZUL)
    add_text(s, val, x + 0.1, 1.5, 2.7, 1.0, size=36, bold=True, color=NARANJA, align=PP_ALIGN.CENTER)
    add_text(s, label, x + 0.1, 2.5, 2.7, 0.6, size=14, color=BLANCO, align=PP_ALIGN.CENTER)

hallazgos = [
    "Taramundi: mayor densidad gastronómica (7,34 locales/1.000 hab)",
    "Oviedo: 353 locales en total — mayor número absoluto de Asturias",
    "Cocina regional: la más frecuente con 461 locales (29% del total con cocina)",
    "Solo 19 locales con diet:vegan relleno — limitación de datos OSM",
    "Concejos rurales del occidente: densidades altas pese a poca población",
]
for i, h in enumerate(hallazgos):
    add_text(s, h, 0.6, 3.4 + i*0.65, 12.5, 0.6, size=16, color=GRIS)

# ── DIAPOSITIVA 8: MODELO ML ─────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, BLANCO)
add_rect(s, 0, 0, 13.33, 1.2, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, TEAL)

add_text(s, "06 · Modelo de clasificación", 0.6, 0.2, 12, 0.8, size=28, bold=True, color=BLANCO)

add_text(s, "Random Forest Classifier", 0.6, 1.35, 6, 0.6, size=22, bold=True, color=AZUL)
add_text(s, "Predice el tipo de local (restaurant / bar / cafe / pub)\na partir de cocina y ciudad",
         0.6, 1.95, 7, 0.8, size=16, color=GRIS)

detalles_m = [
    ("Features", "cocina, ciudad, tiene_telefono, tiene_web, tiene_horario"),
    ("Clases", "restaurant · bar · cafe · pub"),
    ("Entrenamiento", "n_estimators=200, class_weight='balanced'"),
    ("Validación", "Cross-validation 5 folds"),
    ("Reto principal", "Desbalance de clases (restaurant >> resto)"),
]
for i, (k, v) in enumerate(detalles_m):
    add_rect(s, 0.6, 2.9 + i*0.72, 3.2, 0.62, AZUL_CLARO)
    add_rect(s, 3.8, 2.9 + i*0.72, 9.0, 0.62, GRIS_CLARO)
    add_text(s, k, 0.7, 2.95 + i*0.72, 3.0, 0.52, size=15, bold=True, color=BLANCO)
    add_text(s, v, 3.9, 2.95 + i*0.72, 8.8, 0.52, size=15, color=GRIS)

add_text(s, "💡 Solución: class_weight='balanced' para compensar el desbalance",
         0.6, 6.6, 12, 0.5, size=15, bold=True, color=TEAL)

# ── DIAPOSITIVA 9: DASHBOARD ─────────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, BLANCO)
add_rect(s, 0, 0, 13.33, 1.2, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, TEAL)

add_text(s, "07 · Dashboard interactivo", 0.6, 0.2, 12, 0.8, size=28, bold=True, color=BLANCO)
add_text(s, "streamlit run src/dashboard.py  →  http://localhost:8501",
         0.6, 1.35, 12, 0.5, size=16, italic=True, color=AZUL_CLARO)

secciones = [
    ( "KPIs principales", "Total locales, concejos, tipos de cocina, concejo más denso"),
    ("Mapa interactivo", "Filtros por tipo de local y cocina, popup con detalles"),
    ("Comparativa ciudades", "Oviedo vs Gijón vs Avilés por tipo y cocina"),
    ("Cocinas por zona", "Oeste / Centro / Este de Asturias"),
    ("Predicción ML", "Selecciona cocina y ciudad → predice tipo de local"),
    ("Wikidata", "Tabla de platos y productos típicos asturianos"),
]
for i, (ico, titulo, desc) in enumerate(secciones):
    col = i % 2
    row = i // 2
    x = 0.6 + col * 6.3
    y = 2.0 + row * 1.6
    add_rect(s, x, y, 6.0, 1.45, GRIS_CLARO)
    add_rect(s, x, y, 6.0, 0.5, AZUL_CLARO)
    add_text(s, ico + "  " + titulo, x + 0.1, y + 0.05, 5.8, 0.4, size=16, bold=True, color=BLANCO)
    add_text(s, desc, x + 0.15, y + 0.6, 5.7, 0.8, size=14, color=GRIS)

# ── DIAPOSITIVA 10: CONCLUSIONES ─────────────────────────────────────────────
s = blank_slide(prs)
add_rect(s, 0, 0, 13.33, 7.5, AZUL)
add_rect(s, 0, 0, 0.3, 7.5, NARANJA)

add_text(s, "08 · Conclusiones y vías futuras", 0.6, 0.3, 12, 0.8, size=28, bold=True, color=BLANCO)

add_text(s, "Conclusiones", 0.6, 1.3, 6, 0.5, size=20, bold=True, color=NARANJA)
conclusiones = [
    "Pipeline completo ETL → EDA → Modelo → Dashboard",
    "7.568 locales analizados con 4 fuentes de datos abiertas",
    "Taramundi: concejo con mayor densidad gastronómica de Asturias",
    "Cocina regional asturiana: la más frecuente con diferencia",
    "Aplicable en turismo, planificación territorial y análisis económico",
]
for i, c in enumerate(conclusiones):
    add_text(s, "✓  " + c, 0.6, 1.85 + i*0.62, 6.0, 0.55, size=15, color=BLANCO)

add_text(s, "Vías futuras", 7.0, 1.3, 6, 0.5, size=20, bold=True, color=TEAL)
futuras = [
    "Integración con Foursquare (valoraciones)",
    "Despliegue en la nube (Streamlit Cloud)",
    "Automatización del pipeline con cron job",
    "Análisis temporal con datos históricos OSM",
    "Enriquecimiento con datos de HappyCow (vegano)",
]
for i, f in enumerate(futuras):
    add_text(s, "→  " + f, 7.0, 1.85 + i*0.62, 6.0, 0.55, size=15, color=RGBColor(0xBE,0xD9,0xF5))

add_rect(s, 0.6, 6.5, 12.1, 0.6, TEAL)
add_text(s, "github.com/adrimontelope-jpg/GastroAstur",
         0.6, 6.55, 12.1, 0.5, size=16, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# ── GUARDAR ───────────────────────────────────────────────────────────────────
prs.save("presentacion_GastroAstur.pptx")
print("✓ Presentación guardada como presentacion_GastroAstur.pptx")